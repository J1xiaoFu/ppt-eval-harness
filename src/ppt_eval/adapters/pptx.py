"""Safe, deterministic PPTX inspection.

The adapter deliberately keeps Office and model SDKs out of the domain layer.  It
does a bounded ZIP preflight before asking ``python-pptx`` to parse the package.
When that optional dependency is not installed (or encounters an unsupported
shape), a small OOXML reader provides the structural information required by
the deterministic baseline oracles.
"""

from __future__ import annotations

import hashlib
import io
import math
import posixpath
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path, PurePosixPath
from typing import BinaryIO, Iterable, Mapping, Sequence
from xml.etree import ElementTree as ET

DEFAULT_SLIDE_WIDTH = 12_192_000
DEFAULT_SLIDE_HEIGHT = 6_858_000

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
}


class PptxAdapterError(RuntimeError):
    """Base class for PPTX ingestion failures."""


class UnsafePptxError(PptxAdapterError):
    """Raised when the package fails a bounded security preflight."""

    def __init__(self, report: "ZipPreflightReport") -> None:
        self.report = report
        codes = ", ".join(f.code for f in report.blocking_findings)
        super().__init__(f"PPTX package rejected by security preflight: {codes}")


class DependencyUnavailableError(PptxAdapterError):
    """Raised when an explicitly requested parser backend is unavailable."""


@dataclass(frozen=True, slots=True)
class SecurityLimits:
    """Limits are intentionally conservative and can be profile-parameterized."""

    max_archive_bytes: int = 100 * 1024 * 1024
    max_entries: int = 10_000
    max_total_uncompressed_bytes: int = 500 * 1024 * 1024
    max_entry_uncompressed_bytes: int = 100 * 1024 * 1024
    max_compression_ratio: float = 200.0
    max_relationship_scan_bytes: int = 16 * 1024 * 1024

    def __post_init__(self) -> None:
        numeric = (
            self.max_archive_bytes,
            self.max_entries,
            self.max_total_uncompressed_bytes,
            self.max_entry_uncompressed_bytes,
            self.max_relationship_scan_bytes,
        )
        if any(value <= 0 for value in numeric) or self.max_compression_ratio <= 0:
            raise ValueError("all PPTX security limits must be positive")


@dataclass(frozen=True, slots=True)
class SecurityFinding:
    code: str
    severity: str
    message: str
    entry: str | None = None
    payload: Mapping[str, object] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class ZipPreflightReport:
    archive_bytes: int
    entry_count: int
    total_uncompressed_bytes: int
    max_observed_compression_ratio: float
    findings: tuple[SecurityFinding, ...] = ()
    has_macros: bool = False
    has_external_relationships: bool = False

    @property
    def blocking_findings(self) -> tuple[SecurityFinding, ...]:
        return tuple(item for item in self.findings if item.severity == "BLOCK")

    @property
    def is_safe(self) -> bool:
        return not self.blocking_findings


@dataclass(frozen=True, slots=True)
class BoundingBox:
    """Normalized slide coordinates, all values nominally in ``[0, 1]``."""

    x: float
    y: float
    width: float
    height: float

    def as_tuple(self) -> tuple[float, float, float, float]:
        return (self.x, self.y, self.width, self.height)

    @property
    def area(self) -> float:
        return max(0.0, self.width) * max(0.0, self.height)

    @property
    def is_outside_slide(self) -> bool:
        epsilon = 1e-6
        return (
            self.x < -epsilon
            or self.y < -epsilon
            or self.width < 0
            or self.height < 0
            or self.x + self.width > 1 + epsilon
            or self.y + self.height > 1 + epsilon
        )


@dataclass(frozen=True, slots=True)
class SlideObject:
    object_id: str
    name: str
    kind: str
    bbox: BoundingBox
    text: str = ""
    font_sizes_pt: tuple[float, ...] = ()
    font_names: tuple[str, ...] = ()
    editable: bool = True
    hidden: bool = False
    media_sha256: str | None = None
    relationship_target: str | None = None
    metadata: Mapping[str, object] = field(default_factory=dict)

    @property
    def visible_text(self) -> str:
        return "" if self.hidden else self.text.strip()


@dataclass(frozen=True, slots=True)
class ParsedSlide:
    page_number: int
    slide_id: str
    objects: tuple[SlideObject, ...]
    notes_text: str = ""
    hidden: bool = False

    @property
    def visible_objects(self) -> tuple[SlideObject, ...]:
        return tuple(item for item in self.objects if not item.hidden)

    @property
    def visible_text(self) -> str:
        return "\n".join(
            item.visible_text for item in self.objects if item.visible_text
        ).strip()


@dataclass(frozen=True, slots=True)
class ParsedPresentation:
    source_name: str
    source_sha256: str
    width_emu: int
    height_emu: int
    slides: tuple[ParsedSlide, ...]
    media_hashes: tuple[str, ...]
    preflight: ZipPreflightReport
    parser_backend: str
    warnings: tuple[str, ...] = ()
    metadata: Mapping[str, object] = field(default_factory=dict)

    @property
    def slide_count(self) -> int:
        return len(self.slides)

    @property
    def all_visible_text(self) -> str:
        return "\n".join(slide.visible_text for slide in self.slides).strip()

    @property
    def object_count(self) -> int:
        return sum(len(slide.objects) for slide in self.slides)


Source = str | Path | bytes | bytearray | BinaryIO


class PptxAdapter:
    """Parse PPTX packages into vendor-neutral slide objects."""

    def __init__(
        self,
        limits: SecurityLimits | None = None,
        *,
        backend: str = "auto",
    ) -> None:
        if backend not in {"auto", "python-pptx", "ooxml"}:
            raise ValueError("backend must be one of: auto, python-pptx, ooxml")
        self.limits = limits or SecurityLimits()
        self.backend = backend

    def preflight(self, source: Source) -> ZipPreflightReport:
        if isinstance(source, (str, Path)):
            path = Path(source)
            try:
                stat = path.stat()
            except OSError as exc:
                return ZipPreflightReport(
                    archive_bytes=0,
                    entry_count=0,
                    total_uncompressed_bytes=0,
                    max_observed_compression_ratio=0.0,
                    findings=(_finding("unreadable_source", "BLOCK", str(exc)),),
                )
            if not path.is_file():
                return ZipPreflightReport(
                    archive_bytes=0,
                    entry_count=0,
                    total_uncompressed_bytes=0,
                    max_observed_compression_ratio=0.0,
                    findings=(
                        _finding(
                            "unreadable_source",
                            "BLOCK",
                            "PPTX source is not a regular file.",
                        ),
                    ),
                )
            size = stat.st_size
            if size > self.limits.max_archive_bytes:
                return _oversize_report(size, self.limits.max_archive_bytes)
        if isinstance(source, (bytes, bytearray)) and len(source) > self.limits.max_archive_bytes:
            return _oversize_report(len(source), self.limits.max_archive_bytes)
        if hasattr(source, "read") and not isinstance(source, (str, Path, bytes, bytearray)):
            data = source.read(self.limits.max_archive_bytes + 1)
            if not isinstance(data, (bytes, bytearray)):
                raise TypeError("PPTX binary stream must return bytes")
            if len(data) > self.limits.max_archive_bytes:
                return _oversize_report(len(data), self.limits.max_archive_bytes)
            return self.preflight_bytes(bytes(data))
        data, _ = _read_source(source, self.limits.max_archive_bytes)
        return self.preflight_bytes(data)

    def preflight_bytes(self, data: bytes) -> ZipPreflightReport:
        findings: list[SecurityFinding] = []
        if len(data) > self.limits.max_archive_bytes:
            findings.append(
                _finding(
                    "archive_too_large",
                    "BLOCK",
                    "Compressed PPTX exceeds the configured byte limit.",
                    actual=len(data),
                    limit=self.limits.max_archive_bytes,
                )
            )
            return ZipPreflightReport(len(data), 0, 0, 0.0, tuple(findings))

        try:
            archive = zipfile.ZipFile(io.BytesIO(data))
        except (zipfile.BadZipFile, OSError) as exc:
            findings.append(_finding("invalid_zip", "BLOCK", str(exc)))
            return ZipPreflightReport(len(data), 0, 0, 0.0, tuple(findings))

        with archive:
            entries = archive.infolist()
            if len(entries) > self.limits.max_entries:
                findings.append(
                    _finding(
                        "too_many_entries",
                        "BLOCK",
                        "PPTX contains too many ZIP entries.",
                        actual=len(entries),
                        limit=self.limits.max_entries,
                    )
                )

            total = 0
            max_ratio = 0.0
            seen: set[str] = set()
            has_macros = False
            for info in entries:
                normalized = _normalized_entry_name(info.filename)
                lower = normalized.lower()
                if _unsafe_entry_name(info.filename):
                    findings.append(
                        _finding(
                            "unsafe_entry_path",
                            "BLOCK",
                            "ZIP entry attempts path traversal or uses an absolute path.",
                            entry=info.filename,
                        )
                    )
                if lower in seen:
                    findings.append(
                        _finding(
                            "duplicate_entry",
                            "BLOCK",
                            "ZIP contains duplicate case-insensitive entry names.",
                            entry=info.filename,
                        )
                    )
                seen.add(lower)
                if info.flag_bits & 0x1:
                    findings.append(
                        _finding(
                            "encrypted_entry",
                            "BLOCK",
                            "Encrypted ZIP entries are not supported.",
                            entry=info.filename,
                        )
                    )
                total += info.file_size
                if info.file_size > self.limits.max_entry_uncompressed_bytes:
                    findings.append(
                        _finding(
                            "entry_too_large",
                            "BLOCK",
                            "A ZIP entry exceeds the uncompressed byte limit.",
                            entry=info.filename,
                            actual=info.file_size,
                            limit=self.limits.max_entry_uncompressed_bytes,
                        )
                    )
                denominator = max(1, info.compress_size)
                ratio = info.file_size / denominator
                max_ratio = max(max_ratio, ratio)
                if info.file_size > 1_048_576 and ratio > self.limits.max_compression_ratio:
                    findings.append(
                        _finding(
                            "suspicious_compression_ratio",
                            "BLOCK",
                            "A ZIP entry has a suspicious compression ratio.",
                            entry=info.filename,
                            actual=round(ratio, 3),
                            limit=self.limits.max_compression_ratio,
                        )
                    )
                if _is_active_content_name(lower):
                    has_macros = True

            if total > self.limits.max_total_uncompressed_bytes:
                findings.append(
                    _finding(
                        "package_uncompressed_too_large",
                        "BLOCK",
                        "Total uncompressed package size exceeds the configured limit.",
                        actual=total,
                        limit=self.limits.max_total_uncompressed_bytes,
                    )
                )

            required = {"[content_types].xml", "ppt/presentation.xml"}
            for missing in sorted(required - seen):
                findings.append(
                    _finding(
                        "missing_required_part",
                        "BLOCK",
                        f"Required OOXML part is missing: {missing}",
                        entry=missing,
                    )
                )

            if has_macros:
                findings.append(
                    _finding(
                        "active_content_present",
                        "WARN",
                        "Package contains VBA, ActiveX, or embedded object content.",
                    )
                )

            self._scan_xml_directives(archive, entries, findings)
            external_targets = self._scan_external_relationships(archive, entries, findings)
            if external_targets:
                findings.append(
                    _finding(
                        "external_relationships_present",
                        "WARN",
                        "Package contains external relationships; they are never fetched.",
                        count=len(external_targets),
                        sample=external_targets[:10],
                    )
                )

        return ZipPreflightReport(
            archive_bytes=len(data),
            entry_count=len(entries),
            total_uncompressed_bytes=total,
            max_observed_compression_ratio=max_ratio,
            findings=tuple(_deduplicate_findings(findings)),
            has_macros=has_macros,
            has_external_relationships=bool(external_targets),
        )

    def parse(self, source: Source) -> ParsedPresentation:
        data, source_name = _read_source(source, self.limits.max_archive_bytes)
        report = self.preflight_bytes(data)
        if not report.is_safe:
            raise UnsafePptxError(report)

        source_hash = hashlib.sha256(data).hexdigest()
        if self.backend in {"auto", "python-pptx"}:
            try:
                return self._parse_with_python_pptx(
                    data, source_name, source_hash, report
                )
            except ImportError as exc:
                if self.backend == "python-pptx":
                    raise DependencyUnavailableError(
                        "python-pptx is required for the requested parser backend"
                    ) from exc
            except Exception as exc:  # unsupported OOXML should still be assessable
                if self.backend == "python-pptx":
                    raise PptxAdapterError(f"python-pptx failed to parse package: {exc}") from exc
                warning = f"python-pptx fallback: {type(exc).__name__}: {exc}"
                return self._parse_with_ooxml(
                    data, source_name, source_hash, report, warnings=(warning,)
                )

        return self._parse_with_ooxml(data, source_name, source_hash, report)

    def _scan_external_relationships(
        self,
        archive: zipfile.ZipFile,
        entries: Sequence[zipfile.ZipInfo],
        findings: list[SecurityFinding],
    ) -> list[str]:
        scanned = 0
        targets: list[str] = []
        for info in entries:
            if not info.filename.lower().endswith(".rels"):
                continue
            if scanned + info.file_size > self.limits.max_relationship_scan_bytes:
                findings.append(
                    _finding(
                        "relationship_scan_truncated",
                        "WARN",
                        "Relationship scan byte budget was exhausted.",
                    )
                )
                break
            scanned += info.file_size
            try:
                root = ET.fromstring(archive.read(info))
            except (ET.ParseError, RuntimeError, zipfile.BadZipFile) as exc:
                findings.append(
                    _finding(
                        "malformed_relationships",
                        "BLOCK",
                        f"Malformed relationship XML: {exc}",
                        entry=info.filename,
                    )
                )
                continue
            for rel in root:
                if rel.attrib.get("TargetMode", "").lower() == "external":
                    targets.append(rel.attrib.get("Target", ""))
        return targets

    def _scan_xml_directives(
        self,
        archive: zipfile.ZipFile,
        entries: Sequence[zipfile.ZipInfo],
        findings: list[SecurityFinding],
    ) -> None:
        remaining = self.limits.max_relationship_scan_bytes
        for info in entries:
            if not info.filename.lower().endswith((".xml", ".rels")):
                continue
            if remaining <= 0:
                return
            amount = min(info.file_size, 65_536, remaining)
            remaining -= amount
            try:
                with archive.open(info) as stream:
                    prefix = stream.read(amount).upper()
            except (OSError, RuntimeError, zipfile.BadZipFile):
                continue
            if b"<!DOCTYPE" in prefix or b"<!ENTITY" in prefix:
                findings.append(
                    _finding(
                        "unsafe_xml_directive",
                        "BLOCK",
                        "OOXML parts must not contain DTD or entity declarations.",
                        entry=info.filename,
                    )
                )

    def _parse_with_python_pptx(
        self,
        data: bytes,
        source_name: str,
        source_hash: str,
        report: ZipPreflightReport,
    ) -> ParsedPresentation:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(io.BytesIO(data))
        width = int(presentation.slide_width or DEFAULT_SLIDE_WIDTH)
        height = int(presentation.slide_height or DEFAULT_SLIDE_HEIGHT)
        slides: list[ParsedSlide] = []
        media_hashes: set[str] = set()

        for page_number, slide in enumerate(presentation.slides, start=1):
            objects: list[SlideObject] = []
            for shape in slide.shapes:
                shape_type_name = _shape_type_name(shape.shape_type)
                bbox = _bbox_from_emu(
                    _safe_int(getattr(shape, "left", 0)),
                    _safe_int(getattr(shape, "top", 0)),
                    _safe_int(getattr(shape, "width", 0)),
                    _safe_int(getattr(shape, "height", 0)),
                    width,
                    height,
                )
                text = ""
                font_sizes: list[float] = []
                font_names: list[str] = []
                if bool(getattr(shape, "has_text_frame", False)):
                    text_frame = shape.text_frame
                    text = text_frame.text or ""
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            size = getattr(run.font, "size", None)
                            if size is not None:
                                font_sizes.append(round(float(size.pt), 3))
                            name = getattr(run.font, "name", None)
                            if name:
                                font_names.append(str(name))

                if bool(getattr(shape, "has_table", False)):
                    try:
                        table_text = "\n".join(
                            cell.text.strip()
                            for row in shape.table.rows
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        text = "\n".join(part for part in (text, table_text) if part)
                    except (AttributeError, KeyError, ValueError):
                        pass

                chart_values: tuple[str, ...] = ()
                if bool(getattr(shape, "has_chart", False)):
                    chart_values = _python_chart_values(shape)
                    text = "\n".join(part for part in (text, " ".join(chart_values)) if part)

                media_hash = None
                target = None
                image_size_px: tuple[int, int] | None = None
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        media_hash = hashlib.sha256(blob).hexdigest()
                        media_hashes.add(media_hash)
                        target = str(shape.image.filename or "")
                        raw_size = getattr(shape.image, "size", None)
                        if isinstance(raw_size, tuple) and len(raw_size) == 2:
                            image_size_px = (int(raw_size[0]), int(raw_size[1]))
                    except (AttributeError, KeyError, ValueError):
                        pass

                kind = shape_type_name
                if bool(getattr(shape, "has_chart", False)):
                    kind = "chart"
                elif bool(getattr(shape, "has_table", False)):
                    kind = "table"
                editable = kind not in {"picture", "media", "linked_picture"}
                hidden = _python_shape_hidden(shape)
                objects.append(
                    SlideObject(
                        object_id=str(getattr(shape, "shape_id", len(objects) + 1)),
                        name=str(getattr(shape, "name", "")),
                        kind=kind,
                        bbox=bbox,
                        text=text.strip(),
                        font_sizes_pt=tuple(font_sizes),
                        font_names=tuple(dict.fromkeys(font_names)),
                        editable=editable,
                        hidden=hidden,
                        media_sha256=media_hash,
                        relationship_target=target,
                        metadata={
                            "rotation": float(getattr(shape, "rotation", 0.0) or 0.0),
                            "placeholder": bool(getattr(shape, "is_placeholder", False)),
                            "alt_text": _python_shape_alt_text(shape),
                            "crop": _python_picture_crop(shape),
                            "chart_values": chart_values,
                            "image_size_px": image_size_px,
                        },
                    )
                )

            notes_text = ""
            try:
                notes_text = slide.notes_slide.notes_text_frame.text or ""
            except (AttributeError, KeyError, ValueError):
                pass
            slides.append(
                ParsedSlide(
                    page_number=page_number,
                    slide_id=str(getattr(slide, "slide_id", page_number)),
                    objects=tuple(objects),
                    notes_text=notes_text.strip(),
                )
            )

        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            media_hashes.update(_all_media_hashes(archive))

        return ParsedPresentation(
            source_name=source_name,
            source_sha256=source_hash,
            width_emu=width,
            height_emu=height,
            slides=tuple(slides),
            media_hashes=tuple(sorted(media_hashes)),
            preflight=report,
            parser_backend="python-pptx",
            metadata={"format": "pptx"},
        )

    def _parse_with_ooxml(
        self,
        data: bytes,
        source_name: str,
        source_hash: str,
        report: ZipPreflightReport,
        warnings: tuple[str, ...] = (),
    ) -> ParsedPresentation:
        try:
            archive = zipfile.ZipFile(io.BytesIO(data))
            with archive:
                presentation_root = _parse_xml_part(archive, "ppt/presentation.xml")
                size = presentation_root.find("p:sldSz", _NS)
                width = _safe_int(size.attrib.get("cx") if size is not None else None)
                height = _safe_int(size.attrib.get("cy") if size is not None else None)
                width = width or DEFAULT_SLIDE_WIDTH
                height = height or DEFAULT_SLIDE_HEIGHT
                slide_paths = _ordered_slide_paths(archive, presentation_root)
                slides = tuple(
                    _parse_ooxml_slide(archive, path, index, width, height)
                    for index, path in enumerate(slide_paths, start=1)
                )
                media_hashes = tuple(sorted(_all_media_hashes(archive)))
        except (zipfile.BadZipFile, KeyError, ET.ParseError, OSError) as exc:
            raise PptxAdapterError(f"OOXML parser failed: {exc}") from exc

        return ParsedPresentation(
            source_name=source_name,
            source_sha256=source_hash,
            width_emu=width,
            height_emu=height,
            slides=slides,
            media_hashes=media_hashes,
            preflight=report,
            parser_backend="ooxml",
            warnings=warnings,
            metadata={"format": "pptx"},
        )


def _read_source(source: Source, maximum: int) -> tuple[bytes, str]:
    if isinstance(source, (str, Path)):
        path = Path(source)
        try:
            size = path.stat().st_size
        except OSError as exc:
            raise PptxAdapterError(f"cannot read PPTX: {exc}") from exc
        if size > maximum:
            raise PptxAdapterError(
                f"compressed PPTX is {size} bytes; configured limit is {maximum}"
            )
        try:
            return path.read_bytes(), path.name
        except OSError as exc:
            raise PptxAdapterError(f"cannot read PPTX: {exc}") from exc
    if isinstance(source, (bytes, bytearray)):
        data = bytes(source)
        if len(data) > maximum:
            raise PptxAdapterError("compressed PPTX exceeds configured byte limit")
        return data, "memory.pptx"
    if hasattr(source, "read"):
        stream = source
        data = stream.read(maximum + 1)
        if not isinstance(data, (bytes, bytearray)):
            raise TypeError("PPTX binary stream must return bytes")
        if len(data) > maximum:
            raise PptxAdapterError("compressed PPTX exceeds configured byte limit")
        return bytes(data), str(getattr(stream, "name", "stream.pptx"))
    raise TypeError(f"unsupported PPTX source: {type(source)!r}")


def _oversize_report(actual: int, limit: int) -> ZipPreflightReport:
    return ZipPreflightReport(
        archive_bytes=actual,
        entry_count=0,
        total_uncompressed_bytes=0,
        max_observed_compression_ratio=0.0,
        findings=(
            _finding(
                "archive_too_large",
                "BLOCK",
                "Compressed PPTX exceeds the configured byte limit.",
                actual=actual,
                limit=limit,
            ),
        ),
    )


def _finding(
    code: str,
    severity: str,
    message: str,
    *,
    entry: str | None = None,
    **payload: object,
) -> SecurityFinding:
    return SecurityFinding(code, severity, message, entry, payload)


def _deduplicate_findings(findings: Iterable[SecurityFinding]) -> list[SecurityFinding]:
    result: list[SecurityFinding] = []
    seen: set[tuple[str, str | None]] = set()
    for finding in findings:
        key = (finding.code, finding.entry)
        if key not in seen:
            seen.add(key)
            result.append(finding)
    return result


def _normalized_entry_name(name: str) -> str:
    return name.replace("\\", "/").lstrip("/")


def _unsafe_entry_name(name: str) -> bool:
    normalized = name.replace("\\", "/")
    path = PurePosixPath(normalized)
    return (
        normalized.startswith("/")
        or bool(re.match(r"^[A-Za-z]:", normalized))
        or ".." in path.parts
        or "\x00" in normalized
    )


def _is_active_content_name(lower_name: str) -> bool:
    return (
        lower_name.endswith("vbaproject.bin")
        or "/activex/" in f"/{lower_name}"
        or "/embeddings/" in f"/{lower_name}"
        or lower_name.endswith((".exe", ".dll", ".com", ".bat", ".cmd", ".js"))
    )


def _shape_type_name(value: object) -> str:
    name = getattr(value, "name", None)
    if name:
        return str(name).lower()
    text = str(value).split(" ", 1)[0].lower()
    return {
        "13": "picture",
        "14": "placeholder",
        "6": "group",
        "19": "table",
        "3": "chart",
    }.get(text, text or "unknown")


def _safe_int(value: object) -> int:
    try:
        return int(str(value or "0"))
    except (TypeError, ValueError, OverflowError):
        return 0


def _bbox_from_emu(
    x: int, y: int, width: int, height: int, slide_width: int, slide_height: int
) -> BoundingBox:
    sw = max(1, slide_width)
    sh = max(1, slide_height)
    values = (x / sw, y / sh, width / sw, height / sh)
    finite = tuple(value if math.isfinite(value) else 0.0 for value in values)
    return BoundingBox(*finite)


def _python_shape_hidden(shape: object) -> bool:
    try:
        element = getattr(shape, "element")
        c_nv_pr = element.xpath(".//p:cNvPr")[0]
        return str(c_nv_pr.get("hidden", "0")).lower() in {"1", "true"}
    except (AttributeError, IndexError, TypeError):
        return False


def _python_shape_alt_text(shape: object) -> str:
    try:
        element = getattr(shape, "element")
        c_nv_pr = element.xpath(".//p:cNvPr")[0]
        return str(c_nv_pr.get("descr", "") or c_nv_pr.get("title", ""))
    except (AttributeError, IndexError, TypeError):
        return ""


def _python_picture_crop(shape: object) -> tuple[float, float, float, float] | None:
    try:
        values = (
            float(getattr(shape, "crop_left")),
            float(getattr(shape, "crop_top")),
            float(getattr(shape, "crop_right")),
            float(getattr(shape, "crop_bottom")),
        )
    except (AttributeError, TypeError, ValueError):
        return None
    return values


def _python_chart_values(shape: object) -> tuple[str, ...]:
    result: list[str] = []
    try:
        chart = getattr(shape, "chart")
        for series in chart.series:
            name = getattr(series, "name", None)
            if name:
                result.append(str(name))
            values = getattr(series, "values", ())
            for value in values:
                if value is not None:
                    result.append(str(value))
        for plot in chart.plots:
            categories = getattr(plot, "categories", None)
            if categories is not None:
                for category in categories:
                    label = getattr(category, "label", category)
                    if label is not None:
                        result.append(str(label))
    except (AttributeError, KeyError, TypeError, ValueError):
        return tuple(dict.fromkeys(result))
    return tuple(dict.fromkeys(result))


def _parse_xml_part(archive: zipfile.ZipFile, name: str) -> ET.Element:
    data = archive.read(name)
    upper = data[:65_536].upper()
    if b"<!DOCTYPE" in upper or b"<!ENTITY" in upper:
        raise ET.ParseError("DTD and entity declarations are forbidden in OOXML")
    return ET.fromstring(data)


def _relationship_map(archive: zipfile.ZipFile, rel_path: str) -> dict[str, str]:
    try:
        root = _parse_xml_part(archive, rel_path)
    except KeyError:
        return {}
    return {
        item.attrib.get("Id", ""): item.attrib.get("Target", "")
        for item in root
        if item.attrib.get("TargetMode", "").lower() != "external"
    }


def _resolve_part(base_part: str, target: str) -> str:
    base_dir = posixpath.dirname(base_part)
    return posixpath.normpath(posixpath.join(base_dir, target)).lstrip("/")


def _rels_path(part_path: str) -> str:
    directory, filename = posixpath.split(part_path)
    return posixpath.join(directory, "_rels", f"{filename}.rels")


def _ordered_slide_paths(
    archive: zipfile.ZipFile, presentation_root: ET.Element
) -> list[str]:
    relationships = _relationship_map(archive, "ppt/_rels/presentation.xml.rels")
    result: list[str] = []
    for slide_id in presentation_root.findall(".//p:sldId", _NS):
        rel_id = slide_id.attrib.get(f"{{{_NS['r']}}}id", "")
        target = relationships.get(rel_id)
        if target:
            result.append(_resolve_part("ppt/presentation.xml", target))
    if result:
        return result
    names = [
        name
        for name in archive.namelist()
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name, flags=re.IGNORECASE)
    ]
    return sorted(names, key=_slide_part_number)


def _slide_part_number(name: str) -> int:
    match = re.search(r"\d+", name)
    return int(match.group()) if match is not None else 0


def _parse_ooxml_slide(
    archive: zipfile.ZipFile,
    slide_path: str,
    page_number: int,
    slide_width: int,
    slide_height: int,
) -> ParsedSlide:
    root = _parse_xml_part(archive, slide_path)
    rels = _relationship_map(archive, _rels_path(slide_path))
    objects: list[SlideObject] = []
    shape_tags = {
        f"{{{_NS['p']}}}sp": "shape",
        f"{{{_NS['p']}}}pic": "picture",
        f"{{{_NS['p']}}}graphicFrame": "graphic_frame",
        f"{{{_NS['p']}}}grpSp": "group",
        f"{{{_NS['p']}}}cxnSp": "connector",
    }
    for index, element in enumerate(root.iter(), start=1):
        kind = shape_tags.get(element.tag)
        if not kind:
            continue
        c_nv_pr = element.find(".//p:cNvPr", _NS)
        object_id = (
            c_nv_pr.attrib.get("id", str(index)) if c_nv_pr is not None else str(index)
        )
        name = c_nv_pr.attrib.get("name", "") if c_nv_pr is not None else ""
        hidden = (
            c_nv_pr is not None
            and c_nv_pr.attrib.get("hidden", "0").lower() in {"1", "true"}
        )
        xfrm = element.find(".//a:xfrm", _NS)
        if xfrm is None:
            xfrm = element.find(".//p:xfrm", _NS)
        x = y = width = height = 0
        if xfrm is not None:
            offset = xfrm.find("a:off", _NS)
            extent = xfrm.find("a:ext", _NS)
            if offset is not None:
                x = _safe_int(offset.attrib.get("x"))
                y = _safe_int(offset.attrib.get("y"))
            if extent is not None:
                width = _safe_int(extent.attrib.get("cx"))
                height = _safe_int(extent.attrib.get("cy"))
        text = "".join(node.text or "" for node in element.findall(".//a:t", _NS)).strip()
        sizes = tuple(
            sorted(
                {
                    _safe_int(node.attrib.get("sz")) / 100.0
                    for node in element.findall(".//a:rPr", _NS)
                    if _safe_int(node.attrib.get("sz")) > 0
                }
            )
        )
        fonts = tuple(
            dict.fromkeys(
                node.attrib.get("typeface", "")
                for node in element.findall(".//a:latin", _NS)
                if node.attrib.get("typeface")
            )
        )
        chart_values: tuple[str, ...] = ()
        if element.find(".//a:tbl", _NS) is not None:
            kind = "table"
        chart_node = element.find(
            ".//c:chart",
            {**_NS, "c": "http://schemas.openxmlformats.org/drawingml/2006/chart"},
        )
        if chart_node is not None:
            kind = "chart"
        media_hash = None
        relationship_target = None
        embed_node = element.find(".//a:blip", _NS)
        if embed_node is not None:
            rel_id = embed_node.attrib.get(f"{{{_NS['r']}}}embed", "")
            target = rels.get(rel_id)
            if target:
                relationship_target = _resolve_part(slide_path, target)
                try:
                    media_hash = hashlib.sha256(archive.read(relationship_target)).hexdigest()
                except KeyError:
                    pass
        if chart_node is not None:
            rel_id = chart_node.attrib.get(f"{{{_NS['r']}}}id", "")
            target = rels.get(rel_id)
            if target:
                relationship_target = _resolve_part(slide_path, target)
                chart_values = _ooxml_chart_values(archive, relationship_target)
                text = "\n".join(part for part in (text, " ".join(chart_values)) if part)
        objects.append(
            SlideObject(
                object_id=object_id,
                name=name,
                kind=kind,
                bbox=_bbox_from_emu(
                    x, y, width, height, slide_width, slide_height
                ),
                text=text,
                font_sizes_pt=sizes,
                font_names=fonts,
                editable=kind not in {"picture", "media", "linked_picture"},
                hidden=hidden,
                media_sha256=media_hash,
                relationship_target=relationship_target,
                metadata={
                    "alt_text": (
                        c_nv_pr.attrib.get("descr", "") or c_nv_pr.attrib.get("title", "")
                        if c_nv_pr is not None
                        else ""
                    ),
                    "chart_values": chart_values,
                },
            )
        )
    return ParsedSlide(
        page_number=page_number,
        slide_id=posixpath.basename(slide_path),
        objects=tuple(objects),
    )


def _all_media_hashes(archive: zipfile.ZipFile) -> set[str]:
    hashes: set[str] = set()
    for info in archive.infolist():
        if info.filename.lower().startswith("ppt/media/") and not info.is_dir():
            hashes.add(hashlib.sha256(archive.read(info)).hexdigest())
    return hashes


def _ooxml_chart_values(
    archive: zipfile.ZipFile, chart_path: str
) -> tuple[str, ...]:
    try:
        root = _parse_xml_part(archive, chart_path)
    except (KeyError, ET.ParseError):
        return ()
    chart_ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    values = [
        node.text.strip()
        for node in root.findall(".//c:v", chart_ns)
        if node.text and node.text.strip()
    ]
    return tuple(dict.fromkeys(values))
