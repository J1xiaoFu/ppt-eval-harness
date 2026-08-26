from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "demo"


def add_title(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.7), Inches(0.7))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    subtitle_box = slide.shapes.add_textbox(Inches(0.77), Inches(1.32), Inches(11.5), Inches(0.35))
    subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
    subtitle_run.text = subtitle
    subtitle_run.font.name = "Microsoft YaHei"
    subtitle_run.font.size = Pt(11)


def add_bullets(slide, items: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.85), Inches(7.1), Inches(4.6))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(16)
        paragraph.level = 0
    note = slide.shapes.add_textbox(Inches(8.4), Inches(2.0), Inches(3.9), Inches(3.5))
    note_frame = note.text_frame
    note_frame.text = "关键结论"
    note_frame.paragraphs[0].font.size = Pt(16)
    note_frame.paragraphs[0].font.bold = True
    p = note_frame.add_paragraph()
    p.text = items[-1]
    p.font.size = Pt(23)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def generate() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    slides = [
        ("Project Aurora", "项目阶段复盘 · 2026 Q3", ["目标：缩短内容生产周期", "方案：结构化生成与自动质量门禁", "结论：试点周期预计缩短 35%"]),
        ("问题与机会", "从人工抽检转向可解释评测", ["质量标准分散，结果不可复现", "单一模型总分缺少证据和校准", "机会：构建 Oracle 化、可审计的统一 Harness"]),
        ("系统方案", "确定性总控 + 可替换 Oracle", ["基础质量 Oracle 在四场景中强制运行", "专项 Oracle 按证据条件执行", "PDMS 式聚合避免审美掩盖硬错误"]),
        ("验证与下一步", "从元评测走向生产影子运行", ["冻结集验证相关性、稳定性与误放行", "低置信与降级样本转人工复核", "下一步：接入真实编辑 diff，形成数据飞轮"]),
    ]
    for title, subtitle, bullets in slides:
        slide = presentation.slides.add_slide(blank)
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets)

    deck_path = OUT / "aurora_demo.pptx"
    presentation.save(deck_path)

    (OUT / "source.txt").write_text(
        "Project Aurora 目标是缩短内容生产周期。试点预计将周期缩短35%。"
        "系统采用结构化生成、自动质量门禁和人工复核。",
        encoding="utf-8",
    )
    common = {"pptx_path": str(deck_path.resolve()), "audience": "项目评审委员会"}
    cases = {
        "ready_made": {"case_id": "demo-ready-made", "scene": "ready_made", **common},
        "text_to_ppt": {
            "case_id": "demo-text-to-ppt",
            "scene": "text_to_ppt",
            "request": "制作4页中文项目复盘，必须包含目标、系统方案、验证与下一步。",
            **common,
        },
        "project_summary": {
            "case_id": "demo-project-summary",
            "scene": "project_summary",
            "request": "总结项目目标、系统方案和量化收益。",
            "source_materials": [str((OUT / "source.txt").resolve())],
            **common,
        },
        "multimodal": {
            "case_id": "demo-multimodal",
            "scene": "multimodal",
            "request": "使用指定项目素材制作复盘。",
            "assets": [],
            **common,
        },
    }
    for name, payload in cases.items():
        (OUT / f"case_{name}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    print(deck_path)


if __name__ == "__main__":
    generate()
